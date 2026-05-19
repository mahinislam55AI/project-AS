"""
S7-200 PLC Course - Professional PowerPoint Slide Generator
Generates day-wise .pptx files for all 18 modules
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ── Color Palette ──────────────────────────────────────────────────────────────
C_NAVY      = RGBColor(0x00, 0x29, 0x5C)   # Siemens dark blue
C_TEAL      = RGBColor(0x00, 0x99, 0xA1)   # Siemens teal accent
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
C_YELLOW    = RGBColor(0xFF, 0xCC, 0x00)   # Siemens yellow
C_GRAY      = RGBColor(0x64, 0x74, 0x87)
C_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
C_GREEN     = RGBColor(0x00, 0xA8, 0x50)
C_ORANGE    = RGBColor(0xF0, 0x7D, 0x00)
C_RED       = RGBColor(0xCC, 0x00, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Low-level helpers ──────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, l, t, w, h, text, font_size, bold=False,
                color=C_DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return txb

def add_para(tf, text, font_size, bold=False, color=C_DARK_TEXT,
             align=PP_ALIGN.LEFT, bullet=False, indent=0):
    from pptx.oxml.ns import qn
    from lxml import etree
    p   = tf.add_paragraph()
    p.alignment = align
    p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return p


# ── Reusable slide templates ───────────────────────────────────────────────────

def make_title_slide(prs, course_title, subtitle, day_info, module_color=C_NAVY):
    """Full-bleed title slide."""
    slide = blank_slide(prs)
    fill_bg(slide, module_color)

    # top accent bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_YELLOW)

    # big white heading
    add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(1.6),
                course_title, 40, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # teal divider line
    add_rect(slide, Inches(0.6), Inches(2.9), Inches(4), Inches(0.06), C_TEAL)

    # subtitle
    add_textbox(slide, Inches(0.6), Inches(3.1), Inches(11), Inches(0.8),
                subtitle, 22, bold=False, color=C_YELLOW, align=PP_ALIGN.LEFT)

    # day info
    add_textbox(slide, Inches(0.6), Inches(4.0), Inches(11), Inches(0.6),
                day_info, 16, bold=False, color=RGBColor(0xCC,0xDD,0xFF),
                align=PP_ALIGN.LEFT)

    # bottom branding bar
    add_rect(slide, 0, Inches(6.9), SLIDE_W, Inches(0.6), RGBColor(0x00,0x1A,0x40))
    add_textbox(slide, Inches(0.4), Inches(6.92), Inches(7), Inches(0.45),
                "Siemens S7-200 PLC Programming  |  STEP 7-Micro/WIN",
                11, color=C_TEAL, align=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(9), Inches(6.92), Inches(4), Inches(0.45),
                "© S7-200 PLC Course", 11, color=C_GRAY, align=PP_ALIGN.RIGHT)
    return slide


def make_agenda_slide(prs, title, items, day_label, accent=C_TEAL):
    """Agenda / outline slide with numbered items."""
    slide = blank_slide(prs)
    fill_bg(slide, C_LIGHT_BG)

    # left color strip
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, C_NAVY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_YELLOW)

    # header band
    add_rect(slide, Inches(0.35), 0, SLIDE_W, Inches(1.15), C_NAVY)
    add_textbox(slide, Inches(0.55), Inches(0.18), Inches(10), Inches(0.75),
                title, 28, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(10.5), Inches(0.3), Inches(2.6), Inches(0.55),
                day_label, 13, color=C_YELLOW, align=PP_ALIGN.RIGHT)

    # items
    y = Inches(1.35)
    for i, item in enumerate(items):
        num_color = accent
        add_rect(slide, Inches(0.55), y, Inches(0.42), Inches(0.42), num_color)
        add_textbox(slide, Inches(0.56), y - Inches(0.01), Inches(0.42), Inches(0.44),
                    str(i+1), 15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.1), y + Inches(0.02), Inches(11.5), Inches(0.42),
                    item, 16, color=C_DARK_TEXT)
        y += Inches(0.55)

    _footer(slide, day_label)
    return slide


def make_content_slide(prs, title, bullets, day_label, note="", accent=C_TEAL):
    """Standard content slide with bullet points."""
    slide = blank_slide(prs)
    fill_bg(slide, C_WHITE)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_YELLOW)
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, C_NAVY)
    add_rect(slide, Inches(0.35), 0, SLIDE_W, Inches(1.2), C_NAVY)

    add_textbox(slide, Inches(0.55), Inches(0.2), Inches(11.5), Inches(0.8),
                title, 26, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(10.5), Inches(0.3), Inches(2.6), Inches(0.55),
                day_label, 13, color=C_YELLOW, align=PP_ALIGN.RIGHT)

    y = Inches(1.35)
    for b in bullets:
        if b.startswith("##"):          # section header
            add_textbox(slide, Inches(0.55), y, Inches(12.3), Inches(0.38),
                        b[2:].strip(), 15, bold=True, color=accent)
            y += Inches(0.42)
        elif b.startswith("  -"):       # sub-bullet
            add_rect(slide, Inches(1.1), y + Inches(0.15), Inches(0.08), Inches(0.08), accent)
            add_textbox(slide, Inches(1.3), y, Inches(11.5), Inches(0.38),
                        b[3:].strip(), 13, color=C_GRAY)
            y += Inches(0.40)
        else:                           # normal bullet
            add_rect(slide, Inches(0.6), y + Inches(0.13), Inches(0.12), Inches(0.12), accent)
            add_textbox(slide, Inches(0.85), y, Inches(12), Inches(0.4),
                        b.lstrip("- "), 14, color=C_DARK_TEXT)
            y += Inches(0.46)

    if note:
        add_rect(slide, Inches(0.55), Inches(6.5), Inches(12.3), Inches(0.55),
                 RGBColor(0xE8, 0xF4, 0xFF))
        add_textbox(slide, Inches(0.65), Inches(6.52), Inches(12), Inches(0.45),
                    "💡 " + note, 12, color=C_NAVY)

    _footer(slide, day_label)
    return slide


def make_two_col_slide(prs, title, left_title, left_items,
                       right_title, right_items, day_label, accent=C_TEAL):
    """Two-column comparison slide."""
    slide = blank_slide(prs)
    fill_bg(slide, C_WHITE)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_YELLOW)
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, C_NAVY)
    add_rect(slide, Inches(0.35), 0, SLIDE_W, Inches(1.2), C_NAVY)
    add_textbox(slide, Inches(0.55), Inches(0.2), Inches(11.5), Inches(0.8),
                title, 26, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(10.5), Inches(0.3), Inches(2.6), Inches(0.55),
                day_label, 13, color=C_YELLOW, align=PP_ALIGN.RIGHT)

    # Left column
    add_rect(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.45), C_NAVY)
    add_textbox(slide, Inches(0.6), Inches(1.32), Inches(5.6), Inches(0.4),
                left_title, 14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    y = Inches(1.85)
    for item in left_items:
        add_rect(slide, Inches(0.6), y + Inches(0.12), Inches(0.1), Inches(0.1), C_TEAL)
        add_textbox(slide, Inches(0.82), y, Inches(5.3), Inches(0.38), item, 13, color=C_DARK_TEXT)
        y += Inches(0.44)

    # Right column
    add_rect(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(0.45), accent)
    add_textbox(slide, Inches(7.0), Inches(1.32), Inches(5.6), Inches(0.4),
                right_title, 14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    y = Inches(1.85)
    for item in right_items:
        add_rect(slide, Inches(7.0), y + Inches(0.12), Inches(0.1), Inches(0.1), C_NAVY)
        add_textbox(slide, Inches(7.22), y, Inches(5.3), Inches(0.38), item, 13, color=C_DARK_TEXT)
        y += Inches(0.44)

    _footer(slide, day_label)
    return slide


def make_table_slide(prs, title, headers, rows, day_label, accent=C_TEAL):
    """Slide with a data table."""
    slide = blank_slide(prs)
    fill_bg(slide, C_WHITE)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_YELLOW)
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, C_NAVY)
    add_rect(slide, Inches(0.35), 0, SLIDE_W, Inches(1.2), C_NAVY)
    add_textbox(slide, Inches(0.55), Inches(0.2), Inches(11.5), Inches(0.8),
                title, 26, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(10.5), Inches(0.3), Inches(2.6), Inches(0.55),
                day_label, 13, color=C_YELLOW, align=PP_ALIGN.RIGHT)

    col_count = len(headers)
    table_w   = Inches(12.3)
    col_w     = table_w / col_count
    row_h     = Inches(0.45)
    tbl_left  = Inches(0.5)
    tbl_top   = Inches(1.35)

    # Header row
    for ci, h in enumerate(headers):
        add_rect(slide, tbl_left + ci*col_w, tbl_top, col_w, row_h, C_NAVY)
        add_textbox(slide, tbl_left + ci*col_w + Inches(0.05), tbl_top + Inches(0.05),
                    col_w - Inches(0.1), row_h - Inches(0.05),
                    h, 13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        bg = C_LIGHT_BG if ri % 2 == 0 else C_WHITE
        for ci, cell in enumerate(row):
            add_rect(slide, tbl_left + ci*col_w, tbl_top + (ri+1)*row_h,
                     col_w, row_h, bg, C_GRAY)
            add_textbox(slide, tbl_left + ci*col_w + Inches(0.06),
                        tbl_top + (ri+1)*row_h + Inches(0.04),
                        col_w - Inches(0.1), row_h - Inches(0.06),
                        cell, 12, color=C_DARK_TEXT)

    _footer(slide, day_label)
    return slide


def make_practical_slide(prs, title, objective, steps, day_label):
    """Hands-on / practical project slide."""
    slide = blank_slide(prs)
    fill_bg(slide, RGBColor(0xE8, 0xF8, 0xEE))

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_GREEN)
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, C_GREEN)
    add_rect(slide, Inches(0.35), 0, SLIDE_W, Inches(1.2), RGBColor(0x00,0x5C,0x2E))
    add_textbox(slide, Inches(0.55), Inches(0.05), Inches(1.8), Inches(0.5),
                "⚙ PRACTICAL", 13, bold=True, color=C_GREEN)
    add_textbox(slide, Inches(0.55), Inches(0.38), Inches(11.5), Inches(0.72),
                title, 24, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(10.5), Inches(0.3), Inches(2.6), Inches(0.55),
                day_label, 13, color=C_YELLOW, align=PP_ALIGN.RIGHT)

    # Objective box
    add_rect(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.65),
             RGBColor(0x00,0x80,0x45))
    add_textbox(slide, Inches(0.6), Inches(1.38), Inches(12), Inches(0.55),
                "🎯  Objective: " + objective, 14, bold=True, color=C_WHITE)

    # Steps
    y = Inches(2.15)
    for i, step in enumerate(steps):
        add_rect(slide, Inches(0.55), y, Inches(0.38), Inches(0.38),
                 C_GREEN if i % 2 == 0 else RGBColor(0x00,0x70,0x3C))
        add_textbox(slide, Inches(0.56), y, Inches(0.38), Inches(0.38),
                    f"{i+1}", 14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.05), y, Inches(11.7), Inches(0.4),
                    step, 14, color=C_DARK_TEXT)
        y += Inches(0.50)

    _footer(slide, day_label)
    return slide


def make_summary_slide(prs, module_title, key_points, next_module, day_label):
    """End-of-day summary + what's next."""
    slide = blank_slide(prs)
    fill_bg(slide, C_NAVY)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_YELLOW)

    add_textbox(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.7),
                "📋  Session Summary", 28, bold=True, color=C_YELLOW,
                align=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
                module_title, 18, color=C_TEAL)

    add_rect(slide, Inches(0.5), Inches(1.8), Inches(7.6), Inches(4.8),
             RGBColor(0x00,0x1F,0x50))
    add_textbox(slide, Inches(0.65), Inches(1.85), Inches(7.2), Inches(0.45),
                "✅  Key Takeaways", 15, bold=True, color=C_TEAL)
    y = Inches(2.4)
    for kp in key_points:
        add_textbox(slide, Inches(0.75), y, Inches(7.1), Inches(0.4),
                    "▸  " + kp, 13, color=C_WHITE)
        y += Inches(0.44)

    add_rect(slide, Inches(8.4), Inches(1.8), Inches(4.4), Inches(4.8),
             RGBColor(0x00,0x40,0x80))
    add_textbox(slide, Inches(8.55), Inches(1.85), Inches(4.0), Inches(0.45),
                "⏭  Next Session", 15, bold=True, color=C_YELLOW)
    add_textbox(slide, Inches(8.55), Inches(2.45), Inches(4.0), Inches(3.8),
                next_module, 13, color=C_WHITE)

    _footer(slide, day_label)
    return slide


def _footer(slide, day_label):
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), C_NAVY)
    add_textbox(slide, Inches(0.4), Inches(7.17), Inches(6), Inches(0.28),
                "Siemens S7-200 PLC Programming  |  STEP 7-Micro/WIN",
                10, color=C_TEAL)
    add_textbox(slide, Inches(9.5), Inches(7.17), Inches(3.5), Inches(0.28),
                day_label, 10, color=C_YELLOW, align=PP_ALIGN.RIGHT)


def save(prs, filename):
    path = os.path.join(OUTPUT_DIR, filename)
    prs.save(path)
    print(f"  ✅  Saved: {filename}")
