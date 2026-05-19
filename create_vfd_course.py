from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── colour palette ──────────────────────────────────────────────
DELTA_BLUE   = RGBColor(0x00, 0x47, 0xAB)   # deep Delta blue
DELTA_ORANGE = RGBColor(0xFF, 0x6B, 0x00)   # accent orange
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF0, 0xF4, 0xF8)
DARK_GRAY    = RGBColor(0x2D, 0x2D, 0x2D)
MID_BLUE     = RGBColor(0x1A, 0x6E, 0xC8)

# ── helper: solid fill ───────────────────────────────────────────
def solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb

# ── helper: add text box ─────────────────────────────────────────
def tb(slide, text, l, t, w, h, bold=False, size=18,
       color=DARK_GRAY, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    return txb

# ── helper: header band ──────────────────────────────────────────
def header_band(slide, title, subtitle=""):
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0),
                                 prs.slide_width, Inches(1.45))
    solid(bg, DELTA_BLUE)
    bg.line.fill.background()
    tb(slide, title, 0.25, 0.08, 12, 0.72,
       bold=True, size=28, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        tb(slide, subtitle, 0.25, 0.82, 12, 0.55,
           size=16, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.LEFT)

# ── helper: accent bar left ──────────────────────────────────────
def accent_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(1.45),
                                  Inches(0.08), Inches(6.05))
    solid(bar, DELTA_ORANGE)
    bar.line.fill.background()


# ── helper: bullet list ──────────────────────────────────────────
def bullet_list(slide, items, l, t, w, h, size=15, title=None, title_size=17):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = title
        r.font.bold = True; r.font.size = Pt(title_size)
        r.font.color.rgb = DELTA_BLUE
        first = False
    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 1 if item.startswith("  ") else 0
        r = p.add_run()
        prefix = "▶  " if not item.startswith("  ") else "    •  "
        r.text = prefix + item.strip()
        r.font.size = Pt(size)
        r.font.color.rgb = DARK_GRAY
        first = False

# ── helper: two-column bullets ───────────────────────────────────
def two_col(slide, left_items, right_items, left_title="", right_title=""):
    # divider
    div = slide.shapes.add_shape(1, Inches(6.6), Inches(1.6),
                                  Inches(0.04), Inches(5.6))
    solid(div, DELTA_ORANGE); div.line.fill.background()
    bullet_list(slide, left_items,  0.25, 1.55, 6.2, 5.7,
                title=left_title,  size=14)
    bullet_list(slide, right_items, 6.75, 1.55, 6.2, 5.7,
                title=right_title, size=14)

# ── helper: info card ────────────────────────────────────────────
def info_card(slide, heading, lines, l, t, w, h):
    card = slide.shapes.add_shape(1, Inches(l), Inches(t),
                                   Inches(w), Inches(h))
    solid(card, LIGHT_GRAY); card.line.color.rgb = DELTA_BLUE
    card.line.width = Pt(1.2)
    tb(slide, heading, l+0.12, t+0.08, w-0.24, 0.38,
       bold=True, size=14, color=DELTA_BLUE)
    body = "\n".join(lines)
    tb(slide, body, l+0.12, t+0.5, w-0.24, h-0.6,
       size=12, color=DARK_GRAY)

# ════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE / COVER
# ════════════════════════════════════════════════════════════════
def slide_cover():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    solid(bg, DELTA_BLUE); bg.line.fill.background()
    stripe = sl.shapes.add_shape(1, Inches(0), Inches(5.1),
                                  prs.slide_width, Inches(2.4))
    solid(stripe, MID_BLUE); stripe.line.fill.background()
    bar = sl.shapes.add_shape(1, Inches(0), Inches(4.95),
                               prs.slide_width, Inches(0.18))
    solid(bar, DELTA_ORANGE); bar.line.fill.background()
    tb(sl, "Delta VFD-M Series", 0.6, 0.7, 12, 1.1,
       bold=True, size=44, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, "Variable Frequency Drive (VFD)", 0.6, 1.75, 12, 0.65,
       size=26, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
    tb(sl, "Complete Industrial Training Course", 0.6, 2.38, 12, 0.6,
       size=22, color=DELTA_ORANGE, align=PP_ALIGN.CENTER)
    tb(sl, "Sensorless Vector Control Micro Drive  |  AC Motor Speed Control",
       0.6, 3.08, 12, 0.5,
       size=16, color=RGBColor(0xAA,0xCC,0xFF), align=PP_ALIGN.CENTER)
    tb(sl, "Instructor: ________________________     Date: ___________",
       0.6, 5.35, 12, 0.45,
       size=14, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, "Delta Electronics  |  VFD Training Division",
       0.6, 6.0, 12, 0.45,
       size=13, color=RGBColor(0x99,0xBB,0xFF), align=PP_ALIGN.CENTER)
slide_cover()


# ════════════════════════════════════════════════════════════════
#  SLIDE 2 — COURSE OUTLINE (Table of Contents)
# ════════════════════════════════════════════════════════════════
def slide_toc():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg, LIGHT_GRAY); bg.line.fill.background()
    header_band(sl, "Course Outline", "Delta VFD-M Series — Complete Training Program")
    accent_bar(sl)
    modules = [
        ("Module 01", "Introduction to Variable Frequency Drives (VFD)"),
        ("Module 02", "Delta VFD-M Series — Overview & Specifications"),
        ("Module 03", "Safety Guidelines & Precautions"),
        ("Module 04", "Mechanical Installation & Mounting"),
        ("Module 05", "Electrical Wiring — Power & Control Terminals"),
        ("Module 06", "Digital Keypad Operation & Display"),
        ("Module 07", "Parameter Programming (P00 – P09 Groups)"),
        ("Module 08", "Control Modes & Speed Reference"),
        ("Module 09", "Motor Auto-Tuning & Performance Optimization"),
        ("Module 10", "Protective Functions & Fault Diagnosis"),
        ("Module 11", "Industrial Application Examples"),
        ("Module 12", "Maintenance, Troubleshooting & Best Practices"),
    ]
    col1 = modules[:6]
    col2 = modules[6:]
    for i,(mod,title) in enumerate(col1):
        y = 1.65 + i*0.82
        box = sl.shapes.add_shape(1, Inches(0.18), Inches(y),
                                   Inches(6.3), Inches(0.7))
        solid(box, WHITE); box.line.color.rgb = DELTA_BLUE; box.line.width=Pt(0.8)
        tb(sl, mod,   0.22, y+0.06, 1.0, 0.35, bold=True, size=11, color=DELTA_ORANGE)
        tb(sl, title, 1.18, y+0.06, 5.2, 0.55, size=13,   color=DARK_GRAY)
    for i,(mod,title) in enumerate(col2):
        y = 1.65 + i*0.82
        box = sl.shapes.add_shape(1, Inches(6.7), Inches(y),
                                   Inches(6.3), Inches(0.7))
        solid(box, WHITE); box.line.color.rgb = DELTA_BLUE; box.line.width=Pt(0.8)
        tb(sl, mod,   6.74, y+0.06, 1.0, 0.35, bold=True, size=11, color=DELTA_ORANGE)
        tb(sl, title, 7.7,  y+0.06, 5.2, 0.55, size=13,   color=DARK_GRAY)
slide_toc()


# ════════════════════════════════════════════════════════════════
#  MODULE 01 — Introduction to VFD  (2 slides)
# ════════════════════════════════════════════════════════════════
def slide_mod01_a():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 01 — Introduction to Variable Frequency Drives",
                   "What is a VFD? Why do we use it?")
    accent_bar(sl)
    two_col(sl,
        left_title="What is a VFD?",
        left_items=[
            "A VFD (Variable Frequency Drive) is an electronic device that controls the speed of an AC induction motor",
            "It varies the frequency & voltage supplied to the motor",
            "Also known as: Inverter, AC Drive, Frequency Converter",
            "Works on V/F (Volts per Hertz) or Sensorless Vector principle",
            "Input: Fixed AC supply (50/60 Hz)",
            "Output: Variable frequency & variable voltage to motor",
        ],
        right_title="Why Use a VFD?",
        right_items=[
            "Precise motor speed control (0 – max RPM)",
            "Soft start / soft stop → reduces mechanical stress",
            "Energy saving: up to 50% electricity savings on fans & pumps",
            "Reduces inrush current at startup",
            "Overload & short-circuit protection",
            "Forward / Reverse control",
            "Multiple speed presets",
            "Smooth acceleration & deceleration ramps",
        ]
    )
slide_mod01_a()

def slide_mod01_b():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 01 — VFD Working Principle & Types",
                   "How AC drives work internally")
    accent_bar(sl)
    info_card(sl,"VFD Internal Power Flow",
              ["AC Mains → Rectifier (AC→DC) → DC Bus (Capacitor Bank) → IGBT Inverter (DC→Variable AC) → Motor",
               "PWM (Pulse Width Modulation) used to synthesise variable frequency output",
               "Carrier frequency: typically 2 kHz – 15 kHz (P03.00 in VFD-M)"],
              0.18, 1.6, 12.9, 1.15)
    info_card(sl,"Control Methods",
              ["V/F Control: Maintains constant V/Hz ratio — simple, good for fans/pumps",
               "Sensorless Vector Control (SVC): Estimates rotor flux without encoder — high torque at low speed",
               "Delta VFD-M uses SVC as default — superior low-speed torque"],
              0.18, 2.88, 6.1, 2.2)
    info_card(sl,"Common VFD Applications",
              ["Pumps & Water Treatment Plants",
               "HVAC Fans & Air Compressors",
               "Conveyor Belts & Material Handling",
               "CNC Routers, Lathes & Milling Machines",
               "Treadmills & Fitness Equipment",
               "Textile & Packaging Machines"],
              6.55, 2.88, 6.5, 2.2)
    info_card(sl,"Key Terminology",
              ["Frequency (Hz) ↔ Motor Speed (RPM): N = (120 × f) / P",
               "Base Frequency: rated frequency of motor (50 Hz)",
               "Max Output Frequency: can exceed base (field weakening zone)",
               "Acceleration Time: time to ramp from 0 → max speed",
               "Deceleration Time: time to ramp from max speed → 0"],
              0.18, 5.2, 12.9, 1.95)
slide_mod01_b()


# ════════════════════════════════════════════════════════════════
#  MODULE 02 — Delta VFD-M Series Overview  (2 slides)
# ════════════════════════════════════════════════════════════════
def slide_mod02_a():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 02 — Delta VFD-M Series Overview",
                   "Product family, model code, and key specifications")
    accent_bar(sl)
    info_card(sl,"Product Introduction",
              ["Manufacturer: Delta Electronics, Inc. (Taiwan)",
               "Series: VFD-M — General Purpose Sensorless Vector Micro Drive",
               "Technology: Latest microprocessor + IGBT switching",
               "Standard: CE, UL, cUL certified"],
              0.18, 1.6, 6.1, 1.5)
    info_card(sl,"Model Code Breakdown  e.g. VFD007M21A",
              ["VFD   = Variable Frequency Drive product family",
               "007   = 0.75 kW (007 = 0.75kW, 015=1.5kW, 022=2.2kW, 037=3.7kW…)",
               "M     = M Series",
               "2     = Voltage class: 1=115V, 2=230V, 4=460V",
               "1     = Phase: 1=Single phase, 3=Three phase",
               "A     = Version / Enclosure"],
              6.55, 1.6, 6.5, 1.5)
    info_card(sl,"Power & Voltage Ratings",
              ["115V Series:  0.2 kW – 0.75 kW  (0.25 – 1 HP)",
               "230V 1-Phase: 0.4 kW – 2.2 kW   (0.5 – 3 HP)",
               "230V 3-Phase: 0.4 kW – 5.5 kW   (0.5 – 7.5 HP)",
               "460V 3-Phase: 0.75 kW – 5.5 kW  (1 – 7.5 HP)",
               "575V 3-Phase: 0.75 kW – 7.5 kW  (1 – 10 HP)"],
              0.18, 3.22, 6.1, 2.0)
    info_card(sl,"Electrical Specifications",
              ["Output Frequency: 0.1 – 400 Hz",
               "Carrier Frequency: 1 – 15 kHz (adjustable)",
               "Frequency Accuracy: ±0.01 Hz (digital ref)",
               "Overload Capacity: 150% rated current for 60 s",
               "Protection: IP20 / NEMA 1",
               "Operating Temp: -10°C to +40°C (non-condensing)"],
              6.55, 3.22, 6.5, 2.0)
    info_card(sl,"Key Features at a Glance",
              ["Sensorless vector control for high torque at low speed  |  Built-in PID controller",
               "RS-485 Modbus RTU communication  |  Analog input: 0–10V / 4–20mA",
               "Multi-speed presets (up to 15 speeds)  |  Skip frequency bands",
               "DC braking  |  Dynamic braking resistor support  |  Auto-restart after fault"],
              0.18, 5.32, 12.9, 1.8)
slide_mod02_a()

def slide_mod02_b():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 02 — External Components & Accessories",
                   "Parts of the VFD-M unit and optional accessories")
    accent_bar(sl)
    two_col(sl,
        left_title="External Parts of VFD-M Unit",
        left_items=[
            "Digital Keypad (KPV-CE01) — display & programming",
            "Power Terminal Block: R, S, T (Input); U, V, W (Output); B1, B2 (Braking)",
            "Control Terminal Block: Analog & Digital I/O",
            "Grounding Terminal (PE)",
            "Cooling Fan (on larger models)",
            "DIN Rail / Panel Mount Bracket",
            "RS-485 Communication Port (RJ-45)",
            "Charge Indicator LED",
        ],
        right_title="Optional Accessories",
        right_items=[
            "Remote Keypad Extension Cable",
            "Dynamic Braking Resistor (for fast decel)",
            "EMI / RFI Filter (for CE compliance)",
            "AC Line Reactor (input harmonic reduction)",
            "DC Bus Reactor",
            "Zero-Phase Reactor",
            "RS-485 to USB Converter for PC programming",
            "Delta WPLSoft / DIADesigner Software",
        ]
    )
slide_mod02_b()


# ════════════════════════════════════════════════════════════════
#  MODULE 03 — Safety Guidelines  (1 slide)
# ════════════════════════════════════════════════════════════════
def slide_mod03():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 03 — Safety Guidelines & Precautions",
                   "DANGER / WARNING / CAUTION — Read before installation")
    accent_bar(sl)
    # Danger box (red)
    danger = sl.shapes.add_shape(1, Inches(0.18), Inches(1.6),
                                   Inches(12.9), Inches(1.35))
    danger.fill.solid(); danger.fill.fore_color.rgb = RGBColor(0xFF,0xEE,0xEE)
    danger.line.color.rgb = RGBColor(0xCC,0,0); danger.line.width = Pt(1.5)
    tb(sl,"⚠ DANGER",    0.3, 1.65, 2.0, 0.45, bold=True, size=16,
       color=RGBColor(0xCC,0,0))
    tb(sl,"• ALWAYS disconnect AC mains before wiring  "
          "• Wait 10 minutes after power-off for capacitors to discharge  "
          "• Never open the drive while power is ON  "
          "• High voltage present on R/S/T and U/V/W terminals",
       0.3, 2.05, 12.5, 0.8, size=13, color=RGBColor(0x80,0,0))
    # Warning box (orange)
    warn = sl.shapes.add_shape(1, Inches(0.18), Inches(3.08),
                                 Inches(12.9), Inches(1.25))
    warn.fill.solid(); warn.fill.fore_color.rgb = RGBColor(0xFF,0xF5,0xE6)
    warn.line.color.rgb = DELTA_ORANGE; warn.line.width = Pt(1.5)
    tb(sl,"⚠ WARNING",   0.3, 3.13, 2.0, 0.45, bold=True, size=16,
       color=DELTA_ORANGE)
    tb(sl,"• Do NOT connect AC power to U/V/W motor terminals  "
          "• Ground the drive properly (PE terminal)  "
          "• Never touch internal components or PCB  "
          "• Motor cable must be shielded for EMI reduction",
       0.3, 3.5, 12.5, 0.75, size=13, color=RGBColor(0x80,0x40,0))
    # Caution box (blue)
    caut = sl.shapes.add_shape(1, Inches(0.18), Inches(4.45),
                                 Inches(12.9), Inches(1.25))
    caut.fill.solid(); caut.fill.fore_color.rgb = RGBColor(0xE8,0xF0,0xFF)
    caut.line.color.rgb = DELTA_BLUE; caut.line.width = Pt(1.5)
    tb(sl,"ℹ CAUTION",   0.3, 4.5, 2.0, 0.45, bold=True, size=16,
       color=DELTA_BLUE)
    tb(sl,"• Only qualified personnel should install, wire, and program the drive  "
          "• Do not install near heat sources or in high-humidity environments  "
          "• Ensure adequate ventilation — minimum clearance 10 cm on all sides  "
          "• Use copper conductors only; tighten all terminals to specified torque",
       0.3, 4.87, 12.5, 0.75, size=13, color=DELTA_BLUE)
    info_card(sl,"PPE & Tools Required",
              ["Personal Protective Equipment: Insulated gloves, safety glasses, voltage tester",
               "Tools: Insulated screwdrivers, torque wrench, multimeter, clamp meter",
               "Lockout/Tagout (LOTO) procedure must be followed before any maintenance work"],
              0.18, 5.82, 12.9, 1.45)
slide_mod03()


# ════════════════════════════════════════════════════════════════
#  MODULE 04 — Mechanical Installation  (1 slide)
# ════════════════════════════════════════════════════════════════
def slide_mod04():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 04 — Mechanical Installation & Mounting",
                   "Site selection, mounting clearances, and environmental conditions")
    accent_bar(sl)
    two_col(sl,
        left_title="Site & Environmental Requirements",
        left_items=[
            "Ambient Temperature: -10°C to +40°C (derate above 40°C)",
            "Humidity: 90% RH max — non-condensing",
            "Altitude: ≤1000 m (derate above 1000 m)",
            "Vibration: ≤0.6G (20–50 Hz); ≤1G (< 20 Hz)",
            "Storage Temp: -20°C to +60°C",
            "Keep away from: direct sunlight, rain, dust, corrosive gas",
            "Atmospheric Pressure: 86 – 106 kPa",
            "Install in IP54 enclosure for dusty environments",
        ],
        right_title="Mounting Procedure",
        right_items=[
            "Mount VERTICALLY on panel or DIN rail",
            "Minimum clearances: 10 cm top & bottom, 5 cm left & right",
            "Do NOT mount horizontally or upside-down",
            "Use M4 or M5 screws with proper torque",
            "Multiple VFDs: stagger or provide side vents",
            "Ensure cooling air flows bottom to top",
            "Remove the protective film before powering on",
            "Check nameplate voltage matches supply before wiring",
        ]
    )
slide_mod04()

# ════════════════════════════════════════════════════════════════
#  MODULE 05 — Wiring  (2 slides)
# ════════════════════════════════════════════════════════════════
def slide_mod05_a():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 05 — Power Terminal Wiring",
                   "R/S/T input, U/V/W output, braking and grounding")
    accent_bar(sl)
    info_card(sl,"Power Input Terminals (R, S, T)",
              ["R, S, T → 3-phase AC mains input (or R, T for single-phase models)",
               "Always connect through MCCB (Moulded Case Circuit Breaker) or fuse",
               "Never connect motor cable to R/S/T — this will DESTROY the drive",
               "Recommended: add AC line reactor at input for harmonic reduction"],
              0.18, 1.6, 6.1, 1.7)
    info_card(sl,"Power Output Terminals (U, V, W)",
              ["U, V, W → Connect to motor terminals (3-phase)",
               "Do NOT install contactors between drive output and motor",
               "If contactors are needed, interlock with VFD run signal",
               "Use shielded cable; connect shield to PE at both ends"],
              6.55, 1.6, 6.5, 1.7)
    info_card(sl,"Braking & DC Bus Terminals",
              ["B1, B2 → External braking resistor connection",
               "Braking resistor needed for applications with high inertia loads",
               "+1, +2 → DC bus terminals for DC bus reactor (harmonic filter)",
               "Never short-circuit the DC bus terminals"],
              0.18, 3.42, 6.1, 1.75)
    info_card(sl,"Grounding (PE) & Cable Sizing",
              ["Always connect PE (protective earth) to chassis ground",
               "Separate ground conductor required — do NOT share neutral",
               "Cable sizing: follow NEC/IEC based on drive rated current",
               "Tighten terminals to specified torque (see manual table)"],
              6.55, 3.42, 6.5, 1.75)
    info_card(sl,"Wiring Sequence (Safe Practice)",
              ["Step 1: Install MCCB/fuse → Step 2: Wire R/S/T input → Step 3: Wire U/V/W to motor → "
               "Step 4: Connect PE ground → Step 5: Wire control terminals → Step 6: Verify all connections before power-on"],
              0.18, 5.28, 12.9, 1.9)
slide_mod05_a()


def slide_mod05_b():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 05 — Control Terminal Wiring",
                   "Digital inputs, analog inputs, relay outputs, and RS-485")
    accent_bar(sl)
    two_col(sl,
        left_title="Digital Input Terminals",
        left_items=[
            "FWD (MI1) — Forward Run command",
            "REV (MI2) — Reverse Run command",
            "MI3 – MI6 — Multi-function inputs (programmable)",
            "  Common functions: Multi-speed, Jog, External fault, Reset",
            "DCM — Digital common (0V reference)",
            "COM — Common for digital outputs",
            "Logic: NPN (sink) or PNP (source) selectable via jumper",
            "Input voltage: 24VDC internal or external supply",
        ],
        right_title="Analog & Output Terminals",
        right_items=[
            "AVI  — Analog Voltage Input: 0–10V (speed reference)",
            "ACI  — Analog Current Input: 4–20 mA (speed reference)",
            "ACM  — Analog Common",
            "AFM  — Analog Frequency/Current Output (0–10V or 0–20mA)",
            "MO   — Multi-function Open Collector Output",
            "RA, RB, RC — Relay output (250VAC / 30VDC, 1A)",
            "RS-485 port: Modbus RTU, 9600–115200 baud",
            "+10V, GND — for external potentiometer (speed pot)",
        ]
    )
slide_mod05_b()

# ════════════════════════════════════════════════════════════════
#  MODULE 06 — Digital Keypad  (1 slide)
# ════════════════════════════════════════════════════════════════
def slide_mod06():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 06 — Digital Keypad Operation (KPV-CE01)",
                   "Keypad layout, LED display, navigation and status indicators")
    accent_bar(sl)
    two_col(sl,
        left_title="Keypad Buttons & Functions",
        left_items=[
            "▶  RUN    — Start the motor (when keypad control mode)",
            "⏹  STOP   — Stop the motor / Reset fault",
            "▲  UP     — Increase frequency / Scroll parameter up",
            "▼  DOWN   — Decrease frequency / Scroll parameter down",
            "FWD/REV  — Toggle forward/reverse direction",
            "PROG/DATA — Enter parameter programming mode",
            "ENTER    — Confirm parameter value / Enter sub-menu",
            "RESET    — Reset fault condition",
        ],
        right_title="LED Display & Indicators",
        right_items=[
            "5-digit 7-segment LED display",
            "Shows: Output frequency (Hz), Output voltage (V),",
            "  Output current (A), DC bus voltage, RPM",
            "RUN LED — ON when drive is running",
            "FWD LED — ON for forward direction",
            "REV LED — ON for reverse direction",
            "FAULT LED — ON when fault present",
            "Hz / A / V indicator LEDs for unit display",
        ]
    )
slide_mod06()


# ════════════════════════════════════════════════════════════════
#  MODULE 07 — Parameter Programming  (3 slides)
# ════════════════════════════════════════════════════════════════
def slide_mod07_a():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 07 — Parameter Groups Overview (P00 – P09)",
                   "Complete parameter group structure of VFD-M Series")
    accent_bar(sl)
    groups = [
        ("P00","Drive Parameters","Max/Min frequency, base freq, control mode, speed ref source"),
        ("P01","Basic Parameters","Acceleration/Deceleration time (1&2), Jog freq, Jog accel/decel"),
        ("P02","Operation Method","Run/Stop source, direction lock, momentary power loss restart"),
        ("P03","Output Function","Carrier frequency, slip compensation, AVR control, over-voltage stall"),
        ("P04","Input Function","Multi-function input terminal (MI3–MI6) configuration"),
        ("P05","Output Function","Multi-function output (MO, relay RA/RB) configuration"),
        ("P06","Protection","OL1 (motor overload), OC (over-current), OV (over-voltage) settings"),
        ("P07","Motor Parameters","Motor rated voltage, current, frequency, poles, no-load current"),
        ("P08","Special Parameters","DC braking, speed search, PID control enable"),
        ("P09","Communication","RS-485 address, baud rate, parity, Modbus protocol settings"),
    ]
    for i,( grp, name, desc) in enumerate(groups):
        row = sl.shapes.add_shape(1, Inches(0.18), Inches(1.62 + i*0.565),
                                   Inches(12.9), Inches(0.52))
        col = LIGHT_GRAY if i%2==0 else WHITE
        row.fill.solid(); row.fill.fore_color.rgb = col
        row.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); row.line.width=Pt(0.5)
        tb(sl, grp,  0.22, 1.65+i*0.565, 0.7,  0.42, bold=True, size=13, color=DELTA_ORANGE)
        tb(sl, name, 0.95, 1.65+i*0.565, 2.8,  0.42, bold=True, size=13, color=DELTA_BLUE)
        tb(sl, desc, 3.8,  1.65+i*0.565, 9.1,  0.42, size=12,   color=DARK_GRAY)
slide_mod07_a()

def slide_mod07_b():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 07 — Key Parameters: P00 & P01 (Drive & Basic)",
                   "Frequency settings, acceleration/deceleration time programming")
    accent_bar(sl)
    params_p00 = [
        ("P00.00","Max Output Frequency","50.00 Hz","0.10–400.0 Hz","Sets the maximum output frequency"),
        ("P00.01","Motor Base Frequency","50.00 Hz","0.10–400.0 Hz","Motor rated frequency (nameplate)"),
        ("P00.02","Motor Rated Voltage","220V / 380V","Depends on model","Motor nameplate voltage"),
        ("P00.03","Output Freq Limit High","50.00 Hz","0 – P00.00","Upper frequency clamp"),
        ("P00.04","Output Freq Limit Low","0.00 Hz","0 – P00.03","Lower frequency clamp (min speed)"),
        ("P00.05","Control Mode","0","0=V/F, 2=SVC","Select V/F or Sensorless Vector"),
        ("P00.06","Frequency Command Source","0","0=Keypad, 1=AVI, 2=ACI, 3=RS485","Speed reference source"),
        ("P00.07","Run/Stop Command Source","0","0=Keypad, 1=External terminal","Control source selection"),
    ]
    headers = ["Param","Name","Default","Range","Description"]
    widths  = [0.9, 2.5, 1.1, 2.2, 5.7]
    x_pos   = [0.18, 1.1, 3.62, 4.74, 6.96]
    # header row
    hrow = sl.shapes.add_shape(1, Inches(0.18), Inches(1.62),
                                Inches(12.9), Inches(0.42))
    solid(hrow, DELTA_BLUE); hrow.line.fill.background()
    for h,x,w in zip(headers, x_pos, widths):
        tb(sl, h, x, 1.65, w, 0.38, bold=True, size=12, color=WHITE)
    for i,(p,n,d,r,desc) in enumerate(params_p00):
        row = sl.shapes.add_shape(1, Inches(0.18), Inches(2.08+i*0.56),
                                   Inches(12.9), Inches(0.52))
        col = LIGHT_GRAY if i%2==0 else WHITE
        row.fill.solid(); row.fill.fore_color.rgb = col
        row.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); row.line.width=Pt(0.5)
        for val,x,w in zip([p,n,d,r,desc], x_pos, widths):
            bold = (val==p)
            clr  = DELTA_ORANGE if val==p else DARK_GRAY
            tb(sl, val, x, 2.11+i*0.56, w, 0.46, bold=bold, size=11, color=clr)
slide_mod07_b()


def slide_mod07_c():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 07 — Important Parameters: P01, P02, P07",
                   "Accel/Decel, Run source, Motor auto-tune parameters")
    accent_bar(sl)
    info_card(sl,"P01 — Basic Parameters (Accel/Decel)",
              ["P01.00 Accel Time 1:  Default 10s  | Range: 0.1–600s — Time to ramp from 0 to max freq",
               "P01.01 Decel Time 1:  Default 10s  | Range: 0.1–600s — Time to ramp from max to 0",
               "P01.02 Accel Time 2:  Secondary accel time (switched via MI terminal)",
               "P01.03 Decel Time 2:  Secondary decel time",
               "P01.04 Jog Frequency: Default 6 Hz | Used when JOG command activated",
               "P01.05 Jog Accel/Decel: Separate accel/decel time for jog operation"],
              0.18, 1.6, 12.9, 2.15)
    info_card(sl,"P02 — Operation Method",
              ["P02.00 Run/Stop Source: 0=Keypad, 1=Terminal, 2=RS-485",
               "P02.01 Frequency Source: 0=Keypad UP/DN, 1=AVI (0–10V), 2=ACI (4–20mA), 3=RS-485",
               "P02.02 Momentary Power Loss Restart: 0=Disable, 1=Enable (drive restarts after power recovery)",
               "P02.03 Direction Lock: 0=Both, 1=FWD only, 2=REV only"],
              0.18, 3.88, 12.9, 1.85)
    info_card(sl,"P07 — Motor Parameters (for Auto-Tune)",
              ["P07.00 Motor Rated Voltage (V)  — enter motor nameplate value",
               "P07.01 Motor Rated Current (A)  — enter motor nameplate FLA",
               "P07.02 Motor Base Frequency (Hz)— usually 50 Hz",
               "P07.03 Motor Rated RPM          — e.g. 1450 RPM for 4-pole 50 Hz motor",
               "P07.04 Number of Motor Poles    — e.g. 4",
               "P07.05 No-load Current (A)      — measured during auto-tune; improves SVC accuracy"],
              0.18, 5.85, 12.9, 1.85)
slide_mod07_c()

# ════════════════════════════════════════════════════════════════
#  MODULE 08 — Control Modes & Speed Reference  (1 slide)
# ════════════════════════════════════════════════════════════════
def slide_mod08():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 08 — Control Modes & Speed Reference Methods",
                   "Keypad, analog input, multi-speed, RS-485, and PID control")
    accent_bar(sl)
    two_col(sl,
        left_title="Speed Reference Sources",
        left_items=[
            "Keypad Control (P00.06=0): Use ▲/▼ keys to set frequency",
            "Analog Voltage AVI (P00.06=1): 0–10V potentiometer",
            "  0V = 0 Hz,  10V = Max Frequency",
            "Analog Current ACI (P00.06=2): 4–20 mA signal",
            "  4mA = 0 Hz,  20mA = Max Frequency",
            "RS-485 Modbus (P00.06=3): PC or PLC sets frequency via protocol",
            "Multi-Speed Presets (P04 group): Up to 15 fixed speeds",
            "PID Output: Process feedback control (e.g., pressure, flow)",
        ],
        right_title="Run / Stop Command Sources",
        right_items=[
            "Keypad (P02.00=0): RUN / STOP keys on keypad",
            "External Terminal (P02.00=1):",
            "  FWD terminal → motor runs forward",
            "  REV terminal → motor runs reverse",
            "  Both open → motor stops",
            "RS-485 (P02.00=2): Remote Modbus command",
            "Jog Mode: Momentary run at jog frequency (P01.04)",
            "Multi-Function Inputs (MI3–MI6): map to speed, fault reset, etc.",
        ]
    )
slide_mod08()


# ════════════════════════════════════════════════════════════════
#  MODULE 09 — Motor Auto-Tuning  (1 slide)
# ════════════════════════════════════════════════════════════════
def slide_mod09():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 09 — Motor Auto-Tuning & Performance Optimization",
                   "V/F curve, slip compensation, and sensorless vector tuning")
    accent_bar(sl)
    info_card(sl,"Auto-Tuning Procedure (Step-by-Step)",
              ["1. Enter motor nameplate data into P07.00–P07.04",
               "2. Set P05.00 = 12 (Auto-Tuning function) on a Multi-Function output or via keypad",
               "3. Press RUN → Drive will energize motor and measure: stator resistance, leakage inductance",
               "4. Results stored automatically in P07.05 (no-load current) and related registers",
               "5. After tune complete: verify performance by running motor at low speed (5–10 Hz)"],
              0.18, 1.6, 12.9, 2.0)
    info_card(sl,"V/F Curve Settings (P00 group)",
              ["Linear V/F: Constant torque loads (conveyors, compressors) — voltage increases linearly with freq",
               "Square V/F: Variable torque (fans, pumps) — energy efficient, voltage increases as square of freq",
               "Custom V/F: P03.00–P03.03 define up to 4 V/F breakpoints for special motor requirements",
               "Base Frequency (P00.01): Match to motor nameplate (50 Hz in Bangladesh/India)"],
              0.18, 3.72, 6.1, 2.3)
    info_card(sl,"Slip Compensation & DC Braking",
              ["Slip Compensation (P03.07): Compensates for speed drop under load",
               "  Higher value → better speed regulation under varying loads",
               "DC Braking (P08.00–P08.03):",
               "  P08.00 = DC braking current (%)",
               "  P08.01 = DC braking time at start",
               "  P08.02 = DC braking time at stop",
               "  Use for: hoists, cranes, precise stop applications"],
              6.55, 3.72, 6.5, 2.3)
    info_card(sl,"Performance Tips",
              ["Increase carrier frequency (P03.00) for quieter motor operation (reduces audible noise)",
               "Enable AVR (P03.08=1): Auto-Voltage Regulation maintains stable output during input fluctuation",
               "Stall prevention (P06.01): Prevents drive trip during acceleration with heavy loads",
               "Skip frequencies (P04.13–P04.15): Avoid mechanical resonance frequencies"],
              0.18, 6.14, 12.9, 1.6)
slide_mod09()

# ════════════════════════════════════════════════════════════════
#  MODULE 10 — Fault Codes & Protection  (2 slides)
# ════════════════════════════════════════════════════════════════
def slide_mod10_a():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 10 — Protective Functions & Fault Codes (Part 1)",
                   "Over-current, over-voltage, under-voltage, overload faults")
    accent_bar(sl)
    faults_a = [
        ("OC",  "Over-Current",       "Output current exceeded trip level (150–200% rated)","Reduce load, increase accel time, check motor wiring"),
        ("OV",  "Over-Voltage",       "DC bus voltage too high (> ~410V for 230V model)","Add braking resistor, increase decel time"),
        ("UV",  "Under-Voltage",      "DC bus too low — mains supply dropout","Check power supply, check MCCB, check input fuses"),
        ("OL1", "Motor Overload",     "Motor overloaded based on I²t calculation","Reduce load, check motor size, set P06.01 correctly"),
        ("OL2", "Drive Overload",     "Drive itself overloaded — 150% for >60s","Reduce load, upgrade to higher rated drive"),
        ("OH",  "Overheat",           "Heat sink temperature too high (>85°C)","Clean cooling fins, check ambient temp, improve ventilation"),
        ("PHL", "Phase Loss",         "One or more input or output phases missing","Check R/S/T connections, check fuses, check motor cable"),
        ("GFF", "Ground Fault",       "Ground leakage current detected","Check motor insulation, check cable shields"),
    ]
    headers = ["Code","Fault","Cause","Corrective Action"]
    widths  = [0.7, 2.0, 4.1, 5.8]
    x_pos   = [0.18, 0.9, 2.92, 7.04]
    hrow = sl.shapes.add_shape(1, Inches(0.18), Inches(1.62),
                                Inches(12.9), Inches(0.42))
    solid(hrow, DELTA_BLUE); hrow.line.fill.background()
    for h,x,w in zip(headers, x_pos, widths):
        tb(sl, h, x, 1.65, w, 0.38, bold=True, size=13, color=WHITE)
    for i,(code,fault,cause,action) in enumerate(faults_a):
        row = sl.shapes.add_shape(1, Inches(0.18), Inches(2.08+i*0.62),
                                   Inches(12.9), Inches(0.58))
        col = LIGHT_GRAY if i%2==0 else WHITE
        row.fill.solid(); row.fill.fore_color.rgb = col
        row.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); row.line.width=Pt(0.5)
        for val,x,w in zip([code,fault,cause,action], x_pos, widths):
            bold = (val==code)
            clr = RGBColor(0xCC,0,0) if val==code else DARK_GRAY
            tb(sl, val, x, 2.11+i*0.62, w, 0.52, bold=bold, size=11, color=clr)
slide_mod10_a()


def slide_mod10_b():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 10 — Fault Codes (Part 2) & Protection Settings",
                   "Communication faults, hardware faults, and protection parameters")
    accent_bar(sl)
    faults_b = [
        ("EF",  "External Fault",      "Fault signal received on MI terminal (EF function assigned)","Check external device/sensor, clear fault, reset"),
        ("CF1", "Memory Read Error",   "EEPROM internal error","Press RESET; if persists replace keypad/PCB"),
        ("CF3", "CPU Error",           "Internal processor fault","Power cycle drive; contact Delta service"),
        ("LV",  "Low Voltage Warning", "Input voltage dropped below threshold","Check mains voltage; enable auto-restart if needed"),
        ("CE1", "Comm Frame Error",    "RS-485 data frame error","Check baud rate (P09.01), parity (P09.02) settings"),
        ("CE2", "Comm Checksum Error", "Modbus CRC/LRC mismatch","Check cable shielding, cable length < 500m"),
        ("SE",  "Speed Error (SVC)",   "Motor speed feedback error in SVC mode","Re-run auto-tune, check motor parameters"),
    ]
    headers = ["Code","Fault","Cause","Corrective Action"]
    widths  = [0.7, 2.0, 4.1, 5.8]
    x_pos   = [0.18, 0.9, 2.92, 7.04]
    hrow = sl.shapes.add_shape(1, Inches(0.18), Inches(1.62),
                                Inches(12.9), Inches(0.42))
    solid(hrow, DELTA_BLUE); hrow.line.fill.background()
    for h,x,w in zip(headers, x_pos, widths):
        tb(sl, h, x, 1.65, w, 0.38, bold=True, size=13, color=WHITE)
    for i,(code,fault,cause,action) in enumerate(faults_b):
        row = sl.shapes.add_shape(1, Inches(0.18), Inches(2.08+i*0.62),
                                   Inches(12.9), Inches(0.58))
        col = LIGHT_GRAY if i%2==0 else WHITE
        row.fill.solid(); row.fill.fore_color.rgb = col
        row.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); row.line.width=Pt(0.5)
        for val,x,w in zip([code,fault,cause,action], x_pos, widths):
            bold = (val==code)
            clr = RGBColor(0xCC,0,0) if val==code else DARK_GRAY
            tb(sl, val, x, 2.11+i*0.62, w, 0.52, bold=bold, size=11, color=clr)
    info_card(sl,"Fault Reset Methods",
              ["Method 1: Press STOP/RESET key on keypad",
               "Method 2: Toggle EF reset terminal (MI assigned to Reset)",
               "Method 3: Send Modbus reset command via RS-485",
               "Method 4: Power cycle (OFF → wait 5s → ON) — only if fault is cleared"],
              0.18, 6.47, 12.9, 0.88)
slide_mod10_b()


# ════════════════════════════════════════════════════════════════
#  MODULE 11 — Industrial Applications  (2 slides)
# ════════════════════════════════════════════════════════════════
def slide_mod11_a():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 11 — Industrial Application: Pump & Fan Control",
                   "Constant pressure water supply using PID control")
    accent_bar(sl)
    info_card(sl,"Application: Constant Pressure Water Pump System",
              ["Sensor: Pressure transducer → 4–20 mA signal to ACI terminal",
               "Setpoint: Desired pressure entered as P08.10 (PID setpoint)",
               "PID Action: Drive varies motor speed to maintain constant pressure",
               "Benefits: Energy saving, eliminates pressure surges, extends pump life"],
              0.18, 1.6, 12.9, 1.55)
    info_card(sl,"PID Parameter Settings",
              ["P08.10 — PID Setpoint (target value, 0–100%)",
               "P08.11 — PID Feedback Source: 1=ACI (4–20mA pressure transducer)",
               "P08.12 — Proportional Gain (Kp): Higher = faster response",
               "P08.13 — Integral Time (Ti): Lower = more aggressive integration",
               "P08.14 — Derivative Time (Td): Reduces overshoot",
               "P08.15 — PID Output Upper Limit / Lower Limit"],
              0.18, 3.28, 6.1, 2.4)
    info_card(sl,"Fan Speed Control (Energy Saving)",
              ["Use Square V/F curve (P03.06 = 2) for centrifugal fans/pumps",
               "Speed reduction from 100% to 80% → power drops to ~51%",
               "Speed reduction from 100% to 50% → power drops to ~12.5%",
               "P(speed) ∝ Speed³  — Affinity Law for centrifugal loads",
               "ROI on VFD installation: typically < 18 months"],
              6.55, 3.28, 6.5, 2.4)
    info_card(sl,"Typical Parameter Setup: 5.5 kW Water Pump",
              ["P00.01=50Hz, P00.03=50Hz, P00.04=20Hz | P01.00=15s, P01.01=15s | "
               "P02.00=1 (terminal), P02.01=2 (ACI 4-20mA) | "
               "P07.01=11A (motor FLA) | P08.10=60% (setpoint), P08.11=1 (ACI feedback)"],
              0.18, 5.8, 12.9, 1.45)
slide_mod11_a()

def slide_mod11_b():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 11 — Application: CNC / Conveyor / Multi-Speed",
                   "Spindle control, conveyor preset speeds, and RS-485 Modbus")
    accent_bar(sl)
    two_col(sl,
        left_title="CNC Router Spindle Control",
        left_items=[
            "Use AVI (0–10V) from CNC controller to AVI terminal",
            "P00.06 = 1 (AVI frequency source)",
            "P00.02 = 1 (External terminal run/stop)",
            "P00.00 = 400 Hz (max for spindle motor)",
            "P01.00 = 3s, P01.01 = 3s (fast accel/decel)",
            "P03.00 = 12 kHz (high carrier for quiet operation)",
            "P07.01 = spindle motor FLA",
            "Auto-tune P07 parameters for best performance",
        ],
        right_title="Conveyor Multi-Speed Preset",
        right_items=[
            "Up to 15 speed presets via MI terminals",
            "MI3 + MI4 combination selects speed 1–15",
            "P04.05 = Speed 1 (e.g. 10 Hz — slow)",
            "P04.06 = Speed 2 (e.g. 25 Hz — medium)",
            "P04.07 = Speed 3 (e.g. 40 Hz — fast)",
            "P04.08 = Speed 4 (e.g. 50 Hz — full speed)",
            "PLC controls MI terminals for automatic speed selection",
            "Relay output configured to signal 'at speed' status",
        ]
    )
slide_mod11_b()


# ════════════════════════════════════════════════════════════════
#  MODULE 12 — Maintenance & Troubleshooting  (2 slides)
# ════════════════════════════════════════════════════════════════
def slide_mod12_a():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 12 — Preventive Maintenance Schedule",
                   "Daily, monthly, semi-annual, and annual checks")
    accent_bar(sl)
    tasks_daily = [
        "Check for unusual noise or vibration from motor",
        "Monitor output frequency and current on display",
        "Verify cooling fan is running (larger models)",
        "Check ambient temperature is within spec",
    ]
    tasks_monthly = [
        "Clean air vents and heat sink fins (compressed air)",
        "Check all terminal screws for looseness",
        "Inspect power cables for insulation damage",
        "Verify control terminal connections are secure",
    ]
    tasks_semi = [
        "Measure input voltage balance (all 3 phases)",
        "Check motor insulation with megger (>1 MΩ)",
        "Inspect capacitors for bulging or leakage",
        "Test relay output contacts for wear",
    ]
    tasks_annual = [
        "Replace cooling fan if noisy or slow (lifespan ~3 years)",
        "Check electrolytic capacitors (lifespan 5–10 years)",
        "Full parameter backup via keypad or RS-485",
        "Load test drive at full capacity; log current readings",
    ]
    y0 = 1.6
    for title, tasks, l, w in [
        ("Daily Checks", tasks_daily, 0.18, 6.2),
        ("Monthly Checks", tasks_monthly, 6.7, 6.2),
    ]:
        box = sl.shapes.add_shape(1, Inches(l), Inches(y0), Inches(w), Inches(2.5))
        solid(box, LIGHT_GRAY); box.line.color.rgb=DELTA_BLUE; box.line.width=Pt(1)
        tb(sl, title, l+0.12, y0+0.08, w-0.24, 0.4, bold=True, size=14, color=DELTA_BLUE)
        body = "\n".join(["✔  "+t for t in tasks])
        tb(sl, body, l+0.12, y0+0.5, w-0.24, 1.9, size=12, color=DARK_GRAY)
    y0 = 4.22
    for title, tasks, l, w in [
        ("Semi-Annual Checks", tasks_semi, 0.18, 6.2),
        ("Annual Checks", tasks_annual, 6.7, 6.2),
    ]:
        box = sl.shapes.add_shape(1, Inches(l), Inches(y0), Inches(w), Inches(2.5))
        solid(box, LIGHT_GRAY); box.line.color.rgb=DELTA_ORANGE; box.line.width=Pt(1)
        tb(sl, title, l+0.12, y0+0.08, w-0.24, 0.4, bold=True, size=14, color=DELTA_ORANGE)
        body = "\n".join(["✔  "+t for t in tasks])
        tb(sl, body, l+0.12, y0+0.5, w-0.24, 1.9, size=12, color=DARK_GRAY)
slide_mod12_a()

def slide_mod12_b():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg,WHITE); bg.line.fill.background()
    header_band(sl,"Module 12 — Troubleshooting Guide & Best Practices",
                   "Common problems, solutions, and pro tips")
    accent_bar(sl)
    two_col(sl,
        left_title="Common Troubleshooting Scenarios",
        left_items=[
            "Problem: Motor doesn't start",
            "  → Check run source (P02.00), check FWD/REV terminal",
            "  → Check min frequency P00.04 is not too high",
            "Problem: Motor runs but wrong speed",
            "  → Check P00.06 (freq source), verify AVI/ACI signal",
            "Problem: OC fault on start",
            "  → Increase accel time P01.00",
            "  → Check for locked rotor / mechanical jam",
            "Problem: OV fault on stop",
            "  → Increase decel time P01.01",
            "  → Install braking resistor (B1-B2 terminals)",
        ],
        right_title="Best Practices & Pro Tips",
        right_items=[
            "Always back up parameters before changes (note down all Pxx values)",
            "Use shielded cable for motor wiring > 5 metres",
            "Install EMI filter at input for sensitive electronic environments",
            "Never use a contactor on the output side without VFD interlock",
            "Set motor overload P06.01 = motor FLA for proper protection",
            "For PLC integration: use sink/source jumper matching PLC output type",
            "RS-485 Modbus: use 120Ω termination resistor at last device",
            "Always verify drive rating ≥ motor kW before installation",
        ]
    )
slide_mod12_b()


# ════════════════════════════════════════════════════════════════
#  FINAL SLIDE — Summary & Thank You
# ════════════════════════════════════════════════════════════════
def slide_final():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    solid(bg, DELTA_BLUE); bg.line.fill.background()
    stripe = sl.shapes.add_shape(1, Inches(0), Inches(4.7),
                                  prs.slide_width, Inches(2.8))
    solid(stripe, MID_BLUE); stripe.line.fill.background()
    bar = sl.shapes.add_shape(1, Inches(0), Inches(4.56),
                               prs.slide_width, Inches(0.18))
    solid(bar, DELTA_ORANGE); bar.line.fill.background()
    tb(sl,"Course Summary",
       0.5, 0.45, 12.3, 0.7,
       bold=True, size=32, color=WHITE, align=PP_ALIGN.CENTER)
    summary_items = [
        "Module 01: VFD basics, working principle, advantages",
        "Module 02: Delta VFD-M model codes, specs, accessories",
        "Module 03: Safety — DANGER / WARNING / CAUTION guidelines",
        "Module 04: Mechanical installation, clearances, environment",
        "Module 05: Power & control terminal wiring (R/S/T, U/V/W, MI, AVI, ACI)",
        "Module 06: Keypad operation, LED display, navigation",
        "Module 07: Parameter groups P00–P09 — complete programming",
        "Module 08: Speed reference sources & run/stop control modes",
        "Module 09: Auto-tuning, V/F curve, DC braking, optimization",
        "Module 10: Fault codes (OC/OV/UV/OL/OH) & protection settings",
        "Module 11: Applications — pump PID, CNC spindle, multi-speed",
        "Module 12: Maintenance schedule & troubleshooting guide",
    ]
    col1 = summary_items[:6]
    col2 = summary_items[6:]
    for i,item in enumerate(col1):
        tb(sl, "✔  "+item, 0.4, 1.35+i*0.52, 6.2, 0.48, size=12, color=WHITE)
    for i,item in enumerate(col2):
        tb(sl, "✔  "+item, 6.75, 1.35+i*0.52, 6.2, 0.48, size=12, color=WHITE)
    tb(sl,"Thank You for Completing the Delta VFD-M Training Course!",
       0.5, 4.82, 12.3, 0.65,
       bold=True, size=22, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl,"Delta Electronics  |  Industrial Automation Training Division",
       0.5, 5.5, 12.3, 0.45,
       size=15, color=RGBColor(0xAA,0xCC,0xFF), align=PP_ALIGN.CENTER)
    tb(sl,"For technical support: www.delta-ia.com.tw  |  support@deltaww.com",
       0.5, 5.98, 12.3, 0.45,
       size=13, color=RGBColor(0x99,0xBB,0xFF), align=PP_ALIGN.CENTER)
    tb(sl,"© Delta Electronics — All Rights Reserved",
       0.5, 6.55, 12.3, 0.4,
       size=12, color=RGBColor(0x77,0xAA,0xFF), align=PP_ALIGN.CENTER)
slide_final()

# ════════════════════════════════════════════════════════════════
#  SAVE FILE
# ════════════════════════════════════════════════════════════════
output_path = "/projects/sandbox/project-AS/Delta_VFD_M_Series_Course.pptx"
prs.save(output_path)
print(f"✅  Presentation saved → {output_path}")
print(f"    Total slides: {len(prs.slides)}")
